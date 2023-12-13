#!/usr/bin/env python
# coding: utf-8

# In[15]:

from selenium import webdriver
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.common.by import By
from selenium.webdriver.edge.service import Service
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.edge.options import Options
from time import sleep
from datetime import datetime
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import random
import base64
from PIL import Image
import os

#finding the needed folders

script_dir = os.path.abspath(os.path.dirname(__file__))
edgelocationfolder = os.path.join(script_dir,"Edge Webdriver")
edgelocation = os.path.join(edgelocationfolder,"msedgedriver.exe")

logofolder = os.path.join(script_dir,"Logo Folder")
logoslide = os.path.join(script_dir,"Logo Slide Output")

#taking inputs from user

input_companies = input("Enter company names (comma-separated): ")
company_names = [name.strip() for name in input_companies.split(',')]

#disabling console log

op = Options()
op.add_experimental_option('excludeSwitches', ['enable-logging'])

service = Service(executable_path=edgelocation)
#driver = webdriver.Edge(service=service, options=op)

stopmultipledrivers = 0

print("    ")

def searchfunction(companyname):
    
    URL = "https://images.google.com/"
    driver.get(URL)      
    key_input_div = driver.find_element(By.XPATH, ".//textarea[@title='Search']")
    key_input_div.send_keys("'" + companyname + "' full logo transparent wikipedia png")
    searchimage = WebDriverWait(driver, 60).until(EC.element_to_be_clickable((By.XPATH, ".//span[@class='z1asCe MZy1Rb']")))
    driver.execute_script("arguments[0].click();", searchimage)

def findingvalidurl():
       
    firstimage = WebDriverWait(driver, 60).until(EC.presence_of_element_located((By.XPATH, ".//div[@id='islrg']/div/div/a/div/img")))
    src = firstimage.get_attribute('src')
    logo_url.append(src)

def progress_bar(iterable, prefix='', suffix='', length=30, fill='█', print_end='\r'):
    total = len(iterable)
    
    def print_bar(iteration):
        percent = ("{0:.1f}").format(100 * (iteration / float(total)))
        filled_length = int(length * iteration // total)
        bar = fill * filled_length + '-' * (length - filled_length)
        print(f'\r{prefix} |{bar}| {percent}% {suffix}', end=print_end)
    
    for i, item in enumerate(iterable):
        yield item
        print_bar(i + 1)
    
    print()

logo_url = []

def timestampuser():

    username = os.getlogin()
    timestamp = datetime.now()
    timestampformat = timestamp.strftime("%Y-%m-%d %H-%M-%S")
    return timestampformat + " - " + username + ".pptx"

def list_png_images(folder_path):
    png_images = []

    # Ensure the folder path ends with a slash
    folder_path = folder_path.rstrip("/") + "/"

    # List all files in the folder
    files = os.listdir(folder_path)

    # Filter out only the .png files
    png_images = [file for file in files if file.lower().endswith(".png")]

    # Remove ".png" extension from file names
    png_images = [os.path.splitext(file)[0] for file in png_images]

    return png_images

    
print("Finding logos")
for name in progress_bar(range(len(company_names)), prefix='Progress:', suffix='Complete', length=50):
   
    png_images_list = list_png_images(logofolder)
    png_images_list = [item.lower() for item in png_images_list]

    if company_names[name].lower() in png_images_list:
        print("")
        logo_url.append("empty url")
    else:
        
        while(stopmultipledrivers<1):
            driver = webdriver.Edge(service=service, options=op)
            stopmultipledrivers = stopmultipledrivers + 1 
        searchfunction(company_names[name])
        findingvalidurl()

# Disable SSL certificate verification (use with caution)
requests.packages.urllib3.disable_warnings()
logos = dict(zip(company_names, logo_url))

# Create a PowerPoint presentation
prs = Presentation()

# Set the slide size to widescreen 16:9 (10 inches x 5.63 inches)
prs.slide_width = int(10 * 914400)  # 10 inches in EMUs
prs.slide_height = int(5.63 * 914400)  # 5.63 inches in EMUs

slide_width = 10
slide_height = 5.63

# Create a slide with a title and content layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Insert company logos into the slide randomly
# Initialize position variables
row_count = 0
col_count = 0
max_rows = 7
max_cols = 5
left_margin = 0.5  # 1 inch left margin
top_margin = 0.43  # 1 inch top margin
image_width = 0  # 2 inches image width
image_height = 0.3  # 2 inches image height
spacing = 0.3  # 0.25 inches spacing between images
row_spacing = 0.5

left = 0
top = 0

print("Pasting logos in the slide")
for company, url in progress_bar(logos.items(), prefix='Progress:', suffix='Complete', length=50):

    
    if url!="empty url":
        # Base64-encoded image string
        base64_image_string = url
        # Remove the "data:image/png;base64," prefix
        image_data = base64_image_string.split(",")[1]

        # Decode the base64 string to bytes
        image_bytes = base64.b64decode(image_data)

        # Convert the bytes to a PIL Image
        image = Image.open(BytesIO(image_bytes))
        image_file = company + ".png"
        imagestorelocation = os.path.join(logofolder,image_file)
        image.save(imagestorelocation)
    
    sleep(0.75)
    image_file = company + ".png"
    image_path = os.path.join(logofolder, image_file)
    
    try:
        with Image.open(image_path) as img:
            width, height = img.size
            dpi = img.info.get('dpi', (72, 72))  # Default to 72 DPI if not specified
            width_inches = width / dpi[0]
            height_inches = height / dpi[1]
    except Exception as e:
        print(f"Error: {e}")

    # Calculate image position
    
    left = left + left_margin + image_width
    top = top_margin + row_count * (image_height + row_spacing)   
    
    # Maintaing Aspect Ratio
    aspect_ratio = width_inches/height_inches
    image_width = aspect_ratio*image_height
    
    if (left+image_width) < slide_width:    
        pic = slide.shapes.add_picture(image_path, left=Inches(left), top=Inches(top),  height=Inches(image_height), width = Inches(image_width))
    else:
        row_count += 1
        if row_count > (max_rows-1):
            slide = prs.slides.add_slide(slide_layout)
            row_count = 0
            left = 0
            image_width = 0
            top = top_margin + row_count * (image_height + row_spacing)
            left = left + left_margin + image_width
            aspect_ratio = width_inches/height_inches
            image_width = aspect_ratio*image_height
            pic = slide.shapes.add_picture(image_path, left=Inches(left), top=Inches(top),  height=Inches(image_height), width = Inches(image_width))
            
        else:
            left = 0
            image_width = 0
            top = top_margin + row_count * (image_height + row_spacing)
            left = left + left_margin + image_width
            aspect_ratio = width_inches/height_inches
            image_width = aspect_ratio*image_height
            pic = slide.shapes.add_picture(image_path, left=Inches(left), top=Inches(top),  height=Inches(image_height), width = Inches(image_width))
       

slidename = timestampuser()
logoslidelocation = os.path.join(logoslide,slidename)
prs.save(logoslidelocation)
print("    ")
print("Task Done")
print("Ending the script in 5 seconds")
sleep(5)
driver.quit()



# In[ ]:




